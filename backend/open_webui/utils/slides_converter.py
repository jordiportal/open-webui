"""
HTML to PPTX Converter

Converts HTML slides (from Brain artifacts) to PowerPoint format using python-pptx.
"""

import io
import logging
import re
from typing import List, Dict, Any, Optional
from bs4 import BeautifulSoup

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

log = logging.getLogger(__name__)


class SlideData:
    """Parsed data from an HTML slide"""
    def __init__(self):
        self.badge: Optional[str] = None
        self.title: Optional[str] = None
        self.subtitle: Optional[str] = None
        self.content_items: List[str] = []
        self.stats: List[Dict[str, str]] = []  # [{value, label}, ...]
        self.cards: List[Dict[str, str]] = []  # [{num, title, desc}, ...]
        self.conclusion_text: Optional[str] = None


def parse_html_slides(html: str) -> List[SlideData]:
    """Parse HTML content and extract slide data"""
    soup = BeautifulSoup(html, 'html.parser')
    slides_data = []
    
    # Find all slide divs
    slide_divs = soup.find_all('div', class_='slide')
    if not slide_divs:
        # Fallback: try section elements
        slide_divs = soup.find_all('section', class_='slide')
    if not slide_divs:
        # Fallback: try data-slide attribute
        slide_divs = soup.find_all(attrs={'data-slide': True})
    
    for slide_div in slide_divs:
        slide = SlideData()
        
        # Extract badge
        badge = slide_div.find(class_='badge')
        if badge:
            slide.badge = badge.get_text(strip=True)
        
        # Extract title (h1 or h2)
        h1 = slide_div.find('h1')
        h2 = slide_div.find('h2')
        if h1:
            slide.title = h1.get_text(strip=True)
        if h2:
            if not slide.title:
                slide.title = h2.get_text(strip=True)
            else:
                slide.subtitle = h2.get_text(strip=True)
        
        # Extract subtitle from .subtitle class
        subtitle_el = slide_div.find(class_='subtitle')
        if subtitle_el and not slide.subtitle:
            slide.subtitle = subtitle_el.get_text(strip=True)
        
        # Extract list items
        for li in slide_div.find_all('li'):
            text = li.get_text(strip=True)
            if text:
                slide.content_items.append(text)
        
        # Extract stats
        stats_container = slide_div.find(class_='stats')
        if stats_container:
            for stat_div in stats_container.find_all('div', recursive=False):
                value_el = stat_div.find(class_='stat-value')
                label_el = stat_div.find(class_='stat-label')
                if value_el and label_el:
                    slide.stats.append({
                        'value': value_el.get_text(strip=True),
                        'label': label_el.get_text(strip=True)
                    })
        
        # Extract cards
        for card in slide_div.find_all(class_='card'):
            card_data = {}
            num_el = card.find(class_='card-num')
            title_el = card.find(class_='card-title')
            desc_el = card.find(class_='card-desc')
            
            if num_el:
                card_data['num'] = num_el.get_text(strip=True)
            if title_el:
                card_data['title'] = title_el.get_text(strip=True)
            if desc_el:
                card_data['desc'] = desc_el.get_text(strip=True)
            
            if card_data:
                slide.cards.append(card_data)
        
        # Extract conclusion
        conclusion = slide_div.find(class_='conclusion')
        if conclusion:
            slide.conclusion_text = conclusion.get_text(strip=True)
        
        slides_data.append(slide)
    
    return slides_data


def create_pptx(slides_data: List[SlideData], title: str = "Presentation") -> bytes:
    """Create a PPTX file from parsed slide data"""
    prs = Presentation()
    
    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    PURPLE = RGBColor(139, 92, 246)  # #8b5cf6
    DARK_BG = RGBColor(24, 24, 27)   # #18181b
    LIGHT_TEXT = RGBColor(244, 244, 245)  # #f4f4f5
    GRAY_TEXT = RGBColor(161, 161, 170)   # #a1a1aa
    
    for i, slide_data in enumerate(slides_data):
        # Use blank layout for full control
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Set slide background to dark
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG
        
        y_position = Inches(0.5)
        
        # Add badge if present
        if slide_data.badge:
            badge_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), y_position,
                Inches(2.5), Inches(0.4)
            )
            badge_shape.fill.solid()
            badge_shape.fill.fore_color.rgb = RGBColor(88, 28, 135)  # purple-900
            badge_shape.line.fill.background()
            
            tf = badge_shape.text_frame
            tf.paragraphs[0].text = slide_data.badge
            tf.paragraphs[0].font.size = Pt(11)
            tf.paragraphs[0].font.color.rgb = LIGHT_TEXT
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            y_position += Inches(0.6)
        
        # Add title
        if slide_data.title:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), y_position,
                Inches(12), Inches(1)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.title
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = LIGHT_TEXT
            
            y_position += Inches(0.9)
        
        # Add subtitle
        if slide_data.subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.5), y_position,
                Inches(12), Inches(0.5)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data.subtitle
            p.font.size = Pt(18)
            p.font.color.rgb = GRAY_TEXT
            
            y_position += Inches(0.7)
        
        # Add content items (bullet list)
        if slide_data.content_items:
            content_box = slide.shapes.add_textbox(
                Inches(0.5), y_position,
                Inches(12), Inches(3)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for j, item in enumerate(slide_data.content_items):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(18)
                p.font.color.rgb = LIGHT_TEXT
                p.space_after = Pt(12)
            
            y_position += Inches(0.4 * len(slide_data.content_items) + 0.5)
        
        # Add cards in a grid
        if slide_data.cards:
            cards_per_row = 2
            card_width = Inches(5.5)
            card_height = Inches(1.2)
            card_margin = Inches(0.3)
            start_x = Inches(0.5)
            
            for j, card in enumerate(slide_data.cards):
                row = j // cards_per_row
                col = j % cards_per_row
                
                x = start_x + col * (card_width + card_margin)
                y = y_position + row * (card_height + card_margin)
                
                # Card background
                card_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, y, card_width, card_height
                )
                card_shape.fill.solid()
                card_shape.fill.fore_color.rgb = RGBColor(39, 39, 42)  # zinc-800
                card_shape.line.fill.background()
                
                # Card content
                card_text = slide.shapes.add_textbox(
                    x + Inches(0.2), y + Inches(0.15),
                    card_width - Inches(0.4), card_height - Inches(0.3)
                )
                tf = card_text.text_frame
                tf.word_wrap = True
                
                # Number
                if card.get('num'):
                    p = tf.paragraphs[0]
                    p.text = card['num']
                    p.font.size = Pt(14)
                    p.font.bold = True
                    p.font.color.rgb = PURPLE
                
                # Title
                if card.get('title'):
                    p = tf.add_paragraph()
                    p.text = card['title']
                    p.font.size = Pt(14)
                    p.font.bold = True
                    p.font.color.rgb = LIGHT_TEXT
                
                # Description
                if card.get('desc'):
                    p = tf.add_paragraph()
                    p.text = card['desc']
                    p.font.size = Pt(11)
                    p.font.color.rgb = GRAY_TEXT
        
        # Add stats at the bottom
        if slide_data.stats:
            stats_y = Inches(6)
            stat_width = Inches(2)
            start_x = Inches(0.5)
            
            for j, stat in enumerate(slide_data.stats):
                x = start_x + j * (stat_width + Inches(0.5))
                
                stat_box = slide.shapes.add_textbox(
                    x, stats_y, stat_width, Inches(1)
                )
                tf = stat_box.text_frame
                
                # Value
                p = tf.paragraphs[0]
                p.text = stat['value']
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = PURPLE
                
                # Label
                p = tf.add_paragraph()
                p.text = stat['label']
                p.font.size = Pt(10)
                p.font.color.rgb = GRAY_TEXT
        
        # Add conclusion box
        if slide_data.conclusion_text:
            conclusion_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), y_position,
                Inches(12), Inches(1.5)
            )
            conclusion_shape.fill.solid()
            conclusion_shape.fill.fore_color.rgb = RGBColor(88, 28, 135)  # purple-900
            conclusion_shape.line.fill.background()
            
            conclusion_text = slide.shapes.add_textbox(
                Inches(0.7), y_position + Inches(0.2),
                Inches(11.6), Inches(1.1)
            )
            tf = conclusion_text.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.conclusion_text
            p.font.size = Pt(16)
            p.font.color.rgb = LIGHT_TEXT
    
    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


def html_to_pptx(html: str, title: str = "Presentation") -> bytes:
    """
    Convert HTML slides to PPTX format.
    
    Args:
        html: HTML content containing slides (divs with class="slide")
        title: Title of the presentation
    
    Returns:
        PPTX file as bytes
    """
    log.info(f"Converting HTML to PPTX: {title}")
    
    # Parse HTML
    slides_data = parse_html_slides(html)
    
    if not slides_data:
        log.warning("No slides found in HTML, creating empty presentation")
        # Create at least one empty slide
        slides_data = [SlideData()]
        slides_data[0].title = title
    
    # Create PPTX
    pptx_bytes = create_pptx(slides_data, title)
    
    log.info(f"Created PPTX with {len(slides_data)} slides")
    return pptx_bytes
